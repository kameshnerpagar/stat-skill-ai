import io
import pypdf
import docx
import pptx
from typing import Tuple

class DocumentService:
    @staticmethod
    def extract_text(file_bytes: bytes, filename: str) -> Tuple[str, str]:
        """
        Extracts clean text content from PDF, DOCX, PPTX, or TXT binary streams.
        Returns tuple: (extracted_text, file_type)
        """
        ext = filename.split(".")[-1].lower() if "." in filename else ""
        extracted_text = ""

        try:
            if ext == "pdf":
                pdf_reader = pypdf.PdfReader(io.BytesIO(file_bytes))
                text_list = []
                for page in pdf_reader.pages:
                    text = page.extract_text()
                    if text:
                        text_list.append(text)
                extracted_text = "\n\n".join(text_list)
                file_type = "PDF"

            elif ext in ["docx", "doc"]:
                doc = docx.Document(io.BytesIO(file_bytes))
                text_list = [p.text for p in doc.paragraphs if p.text.strip()]
                extracted_text = "\n".join(text_list)
                file_type = "DOCX"

            elif ext in ["pptx", "ppt"]:
                prs = pptx.Presentation(io.BytesIO(file_bytes))
                text_list = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            text_list.append(shape.text)
                extracted_text = "\n".join(text_list)
                file_type = "PPTX"

            else:
                # Default TXT/UTF-8 decoding
                extracted_text = file_bytes.decode("utf-8", errors="ignore")
                file_type = "TXT"

        except Exception as e:
            print(f"[DocumentService] Error extracting file {filename}: {e}")
            extracted_text = file_bytes.decode("utf-8", errors="ignore")
            file_type = ext.upper() or "TXT"

        return extracted_text, file_type
